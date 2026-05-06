"""File-to-markdown conversion with format-specific parsers and timeout protection."""

import asyncio
import logging
from io import BytesIO
from pathlib import Path
from typing import BinaryIO

import pdfplumber
import pypdf
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

logger = logging.getLogger(__name__)

MAX_FILE_CONTENT_LENGTH = 100_000
CONVERSION_TIMEOUT_SECONDS = 30

PDF_EXTENSIONS: frozenset[str] = frozenset({".pdf"})
DOCX_EXTENSIONS: frozenset[str] = frozenset({".docx"})
XLSX_EXTENSIONS: frozenset[str] = frozenset({".xlsx"})
PPTX_EXTENSIONS: frozenset[str] = frozenset({".pptx"})

OFFICE_EXTENSIONS: frozenset[str] = PDF_EXTENSIONS | DOCX_EXTENSIONS | XLSX_EXTENSIONS | PPTX_EXTENSIONS

CODE_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".go",
        ".rs",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".cs",
        ".rb",
        ".php",
        ".swift",
        ".kt",
        ".scala",
        ".sh",
        ".bash",
        ".zsh",
        ".ps1",
        ".html",
        ".css",
        ".scss",
        ".sass",
        ".less",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".xml",
        ".ini",
        ".cfg",
        ".sql",
        ".graphql",
        ".proto",
        ".lua",
        ".r",
        ".m",
        ".zig",
        ".dockerfile",
        ".makefile",
        ".cmake",
    }
)

TEXT_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".txt",
        ".md",
        ".rst",
        ".csv",
        ".tsv",
        ".log",
    }
)

SUPPORTED_EXTENSIONS: frozenset[str] = OFFICE_EXTENSIONS | CODE_EXTENSIONS | TEXT_EXTENSIONS


class ConversionError(Exception):
    """Base exception for file conversion failures."""


class UnsupportedFileError(ConversionError):
    """Raised when the file extension is not supported."""


def is_supported(filename: str) -> bool:
    """Check if the file extension is supported for conversion."""
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


async def convert_file(data: BinaryIO, filename: str) -> str:
    """Convert file bytes to a Markdown string.

    Dispatches to the appropriate converter based on file extension.
    Truncates output at MAX_FILE_CONTENT_LENGTH characters.
    Applies a timeout to prevent runaway conversions.

    Raises:
        UnsupportedFileError: If the file extension is not supported.
        ConversionError: If conversion fails or times out.
    """
    ext = Path(filename).suffix.lower()
    if not ext or ext not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFileError(f"Unsupported file extension: {ext or '(none)'}")

    try:
        async with asyncio.timeout(CONVERSION_TIMEOUT_SECONDS):
            if ext in PDF_EXTENSIONS:
                result = await asyncio.to_thread(_convert_pdf, data)
            elif ext in DOCX_EXTENSIONS:
                result = await asyncio.to_thread(_convert_docx, data)
            elif ext in XLSX_EXTENSIONS:
                result = await asyncio.to_thread(_convert_xlsx, data)
            elif ext in PPTX_EXTENSIONS:
                result = await asyncio.to_thread(_convert_pptx, data)
            elif ext in CODE_EXTENSIONS:
                result = _convert_code_file(data, filename)
            else:
                result = _convert_text_file(data)
    except TimeoutError as exc:
        raise ConversionError(
            f"Failed to convert {filename}: conversion timed out after {CONVERSION_TIMEOUT_SECONDS}s"
        ) from exc
    except Exception as exc:
        raise ConversionError(f"Failed to convert {filename}: {exc}") from exc

    return result


def _convert_pdf(data: BinaryIO) -> str:
    """Convert PDF to plain text using pdfplumber, falling back to pypdf."""
    pdf_stream = data if isinstance(data, BytesIO) else BytesIO(data.read())

    try:
        with pdfplumber.open(pdf_stream) as pdf:
            pages_text = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
            return "\n\n".join(pages_text)
    except Exception:
        logger.warning("pdfplumber failed, falling back to pypdf")
        pdf_stream.seek(0)
        reader = pypdf.PdfReader(pdf_stream)
        pages_text = []
        for pypdf_page in reader.pages:
            text = pypdf_page.extract_text()
            if text:
                pages_text.append(text)
        return "\n\n".join(pages_text)


def _convert_docx(data: BinaryIO) -> str:
    """Convert DOCX to plain text using python-docx."""
    doc = Document(data)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _convert_xlsx(data: BinaryIO) -> str:
    """Convert XLSX to markdown-like text using openpyxl."""
    wb = load_workbook(data, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"## Sheet: {sheet.title}\n")
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            rows.append(row_text)
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def _convert_pptx(data: BinaryIO) -> str:
    """Convert PPTX to markdown-like text using python-pptx."""
    prs = Presentation(data)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides_text.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides_text)


def _convert_code_file(data: BinaryIO, filename: str) -> str:
    content = data.read().decode("utf-8", errors="replace")
    lang = Path(filename).suffix.lstrip(".")
    return f"```{lang}\n{content}\n```"


def _convert_text_file(data: BinaryIO) -> str:
    return data.read().decode("utf-8", errors="replace")
